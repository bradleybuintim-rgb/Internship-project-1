import os
import fitz  # PyMuPDF
from PIL import Image
from docx import Document
from docx.shared import Inches
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
import openpyxl
from pptx import Presentation
from pptx.util import Inches as PptxInches
import tempfile

# ── 1. PDF → WORD ──────────────────────────────────────────────
def pdf_to_word(input_path, output_path):
    doc = fitz.open(input_path)
    word_doc = Document()
    for page in doc:
        text = page.get_text()
        word_doc.add_paragraph(text)
        word_doc.add_page_break()
    word_doc.save(output_path)

# ── 2. PDF → JPG ───────────────────────────────────────────────
def pdf_to_jpg(input_path, output_path):
    doc = fitz.open(input_path)
    page = doc[0]
    mat = fitz.Matrix(2, 2)
    pix = page.get_pixmap(matrix=mat)
    pix.save(output_path)

# ── 3. PDF → PNG ───────────────────────────────────────────────
def pdf_to_png(input_path, output_path):
    doc = fitz.open(input_path)
    page = doc[0]
    mat = fitz.Matrix(2, 2)
    pix = page.get_pixmap(matrix=mat)
    pix.save(output_path)

# ── 4. PDF → EXCEL ─────────────────────────────────────────────
def pdf_to_excel(input_path, output_path):
    doc = fitz.open(input_path)
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "PDF Content"
    row = 1
    for page_num, page in enumerate(doc):
        text = page.get_text()
        ws.cell(row=row, column=1, value=f"--- Page {page_num + 1} ---")
        row += 1
        for line in text.split('\n'):
            if line.strip():
                ws.cell(row=row, column=1, value=line.strip())
                row += 1
        row += 1
    wb.save(output_path)

# ── 5. PDF → PPTX ──────────────────────────────────────────────
def pdf_to_pptx(input_path, output_path):
    doc = fitz.open(input_path)
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]
    for page in doc:
        mat = fitz.Matrix(2, 2)
        pix = page.get_pixmap(matrix=mat)
        tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
        pix.save(tmp.name)
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(tmp.name, 0, 0,
                                  prs.slide_width, prs.slide_height)
        os.unlink(tmp.name)
    prs.save(output_path)

# ── 6. WORD → PDF ──────────────────────────────────────────────
def word_to_pdf(input_path, output_path):
    doc = Document(input_path)
    c = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    y = height - 50
    for para in doc.paragraphs:
        if para.text.strip():
            c.setFont("Helvetica", 11)
            lines = para.text[:100]
            c.drawString(50, y, lines)
            y -= 20
            if y < 50:
                c.showPage()
                y = height - 50
    c.save()

# ── 7. WORD → JPG ──────────────────────────────────────────────
def word_to_jpg(input_path, output_path):
    tmp_pdf = tempfile.NamedTemporaryFile(suffix='.pdf', delete=False)
    word_to_pdf(input_path, tmp_pdf.name)
    pdf_to_jpg(tmp_pdf.name, output_path)
    os.unlink(tmp_pdf.name)

# ── 8. WORD → PNG ──────────────────────────────────────────────
def word_to_png(input_path, output_path):
    tmp_pdf = tempfile.NamedTemporaryFile(suffix='.pdf', delete=False)
    word_to_pdf(input_path, tmp_pdf.name)
    pdf_to_png(tmp_pdf.name, output_path)
    os.unlink(tmp_pdf.name)

# ── 9. JPG → PDF ───────────────────────────────────────────────
def jpg_to_pdf(input_path, output_path):
    image = Image.open(input_path).convert('RGB')
    image.save(output_path, 'PDF')

# ── 10. JPG → PNG ──────────────────────────────────────────────
def jpg_to_png(input_path, output_path):
    image = Image.open(input_path).convert('RGBA')
    image.save(output_path, 'PNG')

# ── 11. PNG → PDF ──────────────────────────────────────────────
def png_to_pdf(input_path, output_path):
    image = Image.open(input_path).convert('RGB')
    image.save(output_path, 'PDF')

# ── 12. PNG → JPG ──────────────────────────────────────────────
def png_to_jpg(input_path, output_path):
    image = Image.open(input_path).convert('RGB')
    image.save(output_path, 'JPEG')

# ── 13. EXCEL → PDF ────────────────────────────────────────────
def excel_to_pdf(input_path, output_path):
    wb = openpyxl.load_workbook(input_path)
    ws = wb.active
    c = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    y = height - 50
    c.setFont("Helvetica-Bold", 13)
    c.drawString(50, y, f"Excel Sheet: {ws.title}")
    y -= 30
    c.setFont("Helvetica", 10)
    for row in ws.iter_rows(values_only=True):
        row_text = "   |   ".join([str(cell) if cell is not None else "" for cell in row])
        if row_text.strip():
            c.drawString(50, y, row_text[:100])
            y -= 18
            if y < 50:
                c.showPage()
                y = height - 50
    c.save()

# ── 14. PPTX → PDF ─────────────────────────────────────────────
def pptx_to_pdf(input_path, output_path):
    prs = Presentation(input_path)
    c = canvas.Canvas(output_path, pagesize=letter)
    width, height = letter
    for slide_num, slide in enumerate(prs.slides):
        y = height - 50
        c.setFont("Helvetica-Bold", 13)
        c.drawString(50, y, f"Slide {slide_num + 1}")
        y -= 30
        c.setFont("Helvetica", 11)
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        c.drawString(50, y, para.text[:100])
                        y -= 20
                        if y < 50:
                            break
        c.showPage()
    c.save()


# ── MASTER CONVERSION FUNCTION ─────────────────────────────────
def convert_file(input_path, conversion_type, output_path):
    conversions = {
        'pdf_to_word':  pdf_to_word,
        'pdf_to_jpg':   pdf_to_jpg,
        'pdf_to_png':   pdf_to_png,
        'pdf_to_excel': pdf_to_excel,
        'pdf_to_pptx':  pdf_to_pptx,
        'word_to_pdf':  word_to_pdf,
        'word_to_jpg':  word_to_jpg,
        'word_to_png':  word_to_png,
        'jpg_to_pdf':   jpg_to_pdf,
        'jpg_to_png':   jpg_to_png,
        'png_to_pdf':   png_to_pdf,
        'png_to_jpg':   png_to_jpg,
        'excel_to_pdf': excel_to_pdf,
        'pptx_to_pdf':  pptx_to_pdf,
    }
    func = conversions.get(conversion_type)
    if func:
        func(input_path, output_path)
    else:
        raise ValueError(f"Unknown conversion type: {conversion_type}")